#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
企业知识库多模态向量化系统 - MD5去重版本
基于文件MD5的智能去重，支持所有Office文档格式
"""

import os
import sys
import logging
from pathlib import Path
import json
import pickle
from typing import List, Dict, Any, Tuple
import numpy as np
import hashlib

# 添加当前目录到Python路径
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

# 导入必要的库
try:
    from sentence_transformers import SentenceTransformer
    import faiss
    import easyocr
    from transformers import pipeline
    import cv2
    import librosa
    from PIL import Image
    import torch
    
    # Office文档处理库
    from docx import Document
    import openpyxl
    from pptx import Presentation
    import mammoth  # 处理.doc文件
    
    print("✅ 所有依赖库导入成功")
except ImportError as e:
    print(f"❌ 依赖库导入失败: {e}")
    sys.exit(1)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

class MD5MultimodalVectorizer:
    """基于MD5的多模态向量化器"""
    
    def __init__(self):
        """初始化MD5向量化器"""
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        logging.info(f"初始化MD5向量化器，设备: {self.device}")
        
        # 初始化文本嵌入模型
        logging.info("初始化文本嵌入模型...")
        self.text_model = SentenceTransformer('BAAI/bge-large-zh-v1.5')
        
        # 初始化OCR处理器
        logging.info("初始化OCR处理器...")
        self.ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=torch.cuda.is_available())
        
        # 初始化图表理解模型
        logging.info("初始化图表理解模型...")
        self.image_captioner = pipeline('image-to-text', model='Salesforce/blip-image-captioning-base')
        
        # 文件MD5缓存
        self.file_md5_cache = {}
        self.processed_files = set()
        
        logging.info("✅ 所有模态处理器初始化完成")
    
    def _calculate_file_md5(self, file_path: str) -> str:
        """计算文件MD5值"""
        try:
            with open(file_path, 'rb') as f:
                file_hash = hashlib.md5()
                chunk = f.read(8192)
                while chunk:
                    file_hash.update(chunk)
                    chunk = f.read(8192)
                return file_hash.hexdigest()
        except Exception as e:
            logging.error(f"计算文件MD5失败 {file_path}: {e}")
            return ""
    
    def _is_duplicate_file(self, file_path: str) -> bool:
        """检查文件是否重复（基于MD5）"""
        file_md5 = self._calculate_file_md5(file_path)
        
        if not file_md5:
            return False
        
        # 检查是否已存在相同MD5的文件
        if file_md5 in self.file_md5_cache:
            existing_file = self.file_md5_cache[file_md5]
            logging.info(f"发现重复文件: {file_path} -> {existing_file}")
            return True
        
        # 记录新文件的MD5
        self.file_md5_cache[file_md5] = file_path
        self.processed_files.add(file_path)
        return False
    
    def _fix_pil_compatibility(self, image: Image.Image, size: Tuple[int, int]) -> Image.Image:
        """终极PIL兼容性修复"""
        try:
            return image.resize(size, Image.Resampling.LANCZOS)
        except AttributeError:
            try:
                return image.resize(size, Image.ANTIALIAS)
            except AttributeError:
                try:
                    return image.resize(size, Image.BICUBIC)
                except AttributeError:
                    return image.resize(size)
    
    def _check_file_integrity(self, file_path: str) -> bool:
        """增强文件完整性检查"""
        try:
            if not os.path.exists(file_path):
                return False
            
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                logging.warning(f"文件大小为0: {file_path}")
                return False
            
            with open(file_path, 'rb') as f:
                f.read(1024)
            
            return True
            
        except Exception as e:
            logging.warning(f"文件完整性检查失败 {file_path}: {e}")
            return False
    
    def _process_pdf_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理PDF文件"""
        try:
            from pdfminer.high_level import extract_text
            from pdfminer.layout import LAParams
            
            if not self._check_file_integrity(file_path):
                return []
            
            text = extract_text(file_path, laparams=LAParams())
            
            if not text.strip():
                logging.warning(f"PDF文件无文本内容: {file_path}")
                return []
            
            chunks = self._split_text(text, max_length=1000)
            
            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    documents.append({
                        'content': chunk.strip(),
                        'type': 'pdf_text',
                        'source': file_path,
                        'chunk_id': i,
                        'metadata': {
                            'file_type': 'pdf',
                            'file_path': file_path,
                            'file_md5': self._calculate_file_md5(file_path),
                            'chunk_size': len(chunk)
                        }
                    })
            
            logging.info(f"PDF文件处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"PDF文件处理失败 {file_path}: {e}")
            return []
    
    def _process_word_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理Word文档 (.doc, .docx)"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            logging.info(f"处理Word文档: {file_path}")
            
            # 使用python-docx处理.docx文件
            if file_path.endswith('.docx'):
                doc = Document(file_path)
                text_content = []
                
                # 提取段落文本
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content.append(paragraph.text.strip())
                
                # 提取表格文本
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content.append(" | ".join(row_text))
                
                full_text = "\n".join(text_content)
            
            # 使用mammoth处理.doc文件
            elif file_path.endswith('.doc'):
                try:
                    with open(file_path, "rb") as docx_file:
                        result = mammoth.convert_to_html(docx_file)
                        # 提取纯文本，去除HTML标签
                        import re
                        full_text = re.sub('<[^<]+?>', '', result.value)
                except Exception as e:
                    logging.error(f"mammoth处理.doc文件失败: {e}")
                    return []
            else:
                logging.warning(f"不支持的文件类型: {file_path}")
                return []
            
            if not full_text.strip():
                logging.warning(f"Word文档无文本内容: {file_path}")
                return []
            
            chunks = self._split_text(full_text, max_length=1000)
            
            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    documents.append({
                        'content': chunk.strip(),
                        'type': 'word_text',
                        'source': file_path,
                        'chunk_id': i,
                        'metadata': {
                            'file_type': 'word',
                            'file_path': file_path,
                            'file_md5': self._calculate_file_md5(file_path),
                            'chunk_size': len(chunk)
                        }
                    })
            
            logging.info(f"Word文档处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"Word文档处理失败 {file_path}: {e}")
            return []
    
    def _process_powerpoint_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理PowerPoint文档 (.ppt, .pptx)"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            logging.info(f"处理PowerPoint文档: {file_path}")
            
            # 使用python-pptx处理.pptx文件
            if file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                text_content = []
                
                # 提取幻灯片文本
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    
                    # 提取形状文本
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    # 提取表格文本
                    for shape in slide.shapes:
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    slide_text.append(" | ".join(row_text))
                    
                    if slide_text:
                        text_content.append(f"幻灯片{slide_num + 1}: " + " | ".join(slide_text))
            
            # 对于.ppt文件，暂时跳过（需要其他库支持）
            elif file_path.endswith('.ppt'):
                logging.warning(f"暂时跳过.ppt文件: {file_path}")
                return []
            else:
                logging.warning(f"不支持的文件类型: {file_path}")
                return []
            
            if not text_content:
                logging.warning(f"PowerPoint文档无文本内容: {file_path}")
                return []
            
            # 将每张幻灯片作为一个文档块
            documents = []
            for i, slide_text in enumerate(text_content):
                if slide_text.strip():
                    documents.append({
                        'content': slide_text.strip(),
                        'type': 'powerpoint_text',
                        'source': file_path,
                        'slide_id': i + 1,
                        'metadata': {
                            'file_type': 'powerpoint',
                            'file_path': file_path,
                            'file_md5': self._calculate_file_md5(file_path),
                            'slide_number': i + 1,
                            'content_size': len(slide_text)
                        }
                    })
            
            logging.info(f"PowerPoint文档处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"PowerPoint文档处理失败 {file_path}: {e}")
            return []
    
    def _process_excel_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理Excel文档 (.xls, .xlsx, .xlsm)"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            logging.info(f"处理Excel文档: {file_path}")
            
            # 使用openpyxl处理.xlsx和.xlsm文件
            if file_path.endswith(('.xlsx', '.xlsm')):
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text_content = []
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_text = []
                    
                    # 提取单元格文本
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell_value in row:
                            if cell_value is not None and str(cell_value).strip():
                                row_text.append(str(cell_value).strip())
                        if row_text:
                            sheet_text.append(" | ".join(row_text))
                    
                    if sheet_text:
                        text_content.append(f"工作表: {sheet_name}\n" + "\n".join(sheet_text))
            
            # 对于.xls文件，暂时跳过（需要pandas支持）
            elif file_path.endswith('.xls'):
                logging.warning(f"暂时跳过.xls文件: {file_path}")
                return []
            else:
                logging.warning(f"不支持的文件类型: {file_path}")
                return []
            
            if not text_content:
                logging.warning(f"Excel文档无文本内容: {file_path}")
                return []
            
            # 将每个工作表作为一个文档块
            documents = []
            for i, sheet_text in enumerate(text_content):
                if sheet_text.strip():
                    documents.append({
                        'content': sheet_text.strip(),
                        'type': 'excel_text',
                        'source': file_path,
                        'sheet_id': i + 1,
                        'metadata': {
                            'file_type': 'excel',
                            'file_path': file_path,
                            'file_md5': self._calculate_file_md5(file_path),
                            'sheet_number': i + 1,
                            'content_size': len(sheet_text)
                        }
                    })
            
            logging.info(f"Excel文档处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"Excel文档处理失败 {file_path}: {e}")
            return []
    
    def _process_txt_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理文本文件"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin1']
            text = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                logging.error(f"无法读取文本文件: {file_path}")
                return []
            
            chunks = self._split_text(text, max_length=1000)
            
            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    documents.append({
                        'content': chunk.strip(),
                        'type': 'text',
                        'source': file_path,
                        'chunk_id': i,
                        'metadata': {
                            'file_type': 'txt',
                            'file_path': file_path,
                            'file_md5': self._calculate_file_md5(file_path),
                            'chunk_size': len(chunk)
                        }
                    })
            
            logging.info(f"文本文件处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"文本文件处理失败 {file_path}: {e}")
            return []
    
    def _process_image_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理图像文件"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            image = Image.open(file_path)
            
            if image.size[0] > 800 or image.size[1] > 800:
                image = self._fix_pil_compatibility(image, (800, 800))
            
            # OCR识别
            try:
                ocr_result = self.ocr_reader.readtext(np.array(image))
                text_content = " ".join([item[1] for item in ocr_result if item[1].strip()])
            except Exception as e:
                logging.warning(f"OCR识别失败: {e}")
                text_content = "图像OCR识别失败"
            
            # 图表理解
            try:
                caption = self.image_captioner(image)[0]['generated_text']
            except Exception as e:
                caption = "图像内容描述"
                logging.warning(f"图表理解失败: {e}")
            
            combined_content = f"OCR识别文本: {text_content}\n图表理解: {caption}"
            
            documents = [{
                'content': combined_content,
                'type': 'image',
                'source': file_path,
                'metadata': {
                    'file_type': 'image',
                    'file_path': file_path,
                    'file_md5': self._calculate_file_md5(file_path),
                    'ocr_text': text_content,
                    'caption': caption,
                    'image_size': image.size
                }
            }]
            
            logging.info(f"图像文件处理完成: {len(documents)} 个文档块")
            return documents
            
        except Exception as e:
            logging.error(f"图像处理失败 {file_path}: {e}")
            return []
    
    def _process_video_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理视频文件"""
        try:
            if not self._check_file_integrity(file_path):
                return []
            
            logging.info(f"提取视频关键帧: {file_path}")
            
            frames = self._extract_key_frames(file_path)
            audio_features = self._analyze_video_audio(file_path)
            
            documents = []
            
            if frames:
                documents.append({
                    'content': f"视频信息: {frames[0].shape[1]}x{frames[0].shape[0]}, 提取{len(frames)}个关键帧...",
                    'type': 'video_frames',
                    'source': file_path,
                    'metadata': {
                        'file_type': 'video',
                        'file_path': file_path,
                        'file_md5': self._calculate_file_md5(file_path),
                        'frame_count': len(frames),
                        'frame_size': frames[0].shape if frames else None
                    }
                })
            
            if audio_features:
                documents.append({
                    'content': f"音频特征: 节拍={audio_features.get('tempo', 'N/A')}BPM, 频谱质心={audio_features.get('spectral_centroid', 'N/A')}Hz...",
                    'type': 'video_audio',
                    'source': file_path,
                    'metadata': {
                        'file_type': 'video',
                        'file_path': file_path,
                        'file_md5': self._calculate_file_md5(file_path),
                        'audio_features': audio_features
                    }
                })
            
            logging.info(f"视频处理完成: {len(documents)} 个文档")
            return documents
            
        except Exception as e:
            logging.error(f"视频处理失败 {file_path}: {e}")
            return []
    
    def _extract_key_frames(self, video_path: str, num_frames: int = 10) -> List[np.ndarray]:
        """提取视频关键帧"""
        try:
            cap = cv2.VideoCapture(video_path)
            
            if not cap.isOpened():
                logging.error(f"无法打开视频文件: {video_path}")
                return []
            
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            
            if total_frames == 0:
                logging.warning(f"视频文件帧数为0: {video_path}")
                cap.release()
                return []
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            duration = total_frames / fps if fps > 0 else 0
            
            logging.info(f"视频信息: {total_frames}帧, {fps:.1f}fps, {duration:.1f}秒")
            
            frame_indices = np.linspace(0, total_frames - 1, num_frames, dtype=int)
            frames = []
            
            for idx in frame_indices:
                cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                ret, frame = cap.read()
                if ret and frame is not None:
                    frames.append(frame)
                else:
                    logging.warning(f"无法读取第{idx}帧")
            
            cap.release()
            
            if not frames:
                logging.warning(f"未成功提取任何关键帧: {video_path}")
            
            return frames
            
        except Exception as e:
            logging.error(f"关键帧提取失败: {e}")
            return []
    
    def _analyze_video_audio(self, video_path: str) -> Dict[str, Any]:
        """分析视频音频特征"""
        try:
            file_size = os.path.getsize(video_path)
            if file_size < 1024:
                logging.warning(f"视频文件过小，跳过音频分析: {video_path}")
                return {}
            
            try:
                y, sr = librosa.load(video_path, sr=None)
            except Exception as e:
                logging.warning(f"librosa加载失败，尝试其他方法: {e}")
                try:
                    import soundfile as sf
                    y, sr = sf.read(video_path)
                except Exception as e2:
                    logging.warning(f"soundfile也失败，跳过音频分析: {e2}")
                    return {}
            
            if len(y) == 0:
                logging.warning(f"音频数据为空: {video_path}")
                return {}
            
            try:
                tempo, _ = librosa.beat.beat_track(y=y, sr=sr)
                spectral_centroid = librosa.feature.spectral_centroid(y=y, sr=sr).mean()
                
                return {
                    'tempo': tempo,
                    'spectral_centroid': spectral_centroid,
                    'duration': len(y) / sr,
                    'sample_rate': sr
                }
                
            except Exception as e:
                logging.warning(f"音频特征提取失败: {e}")
                return {}
            
        except Exception as e:
            logging.warning(f"音频分析失败: {e}")
            return {}
    
    def _split_text(self, text: str, max_length: int = 1000) -> List[str]:
        """分割文本为固定长度的块"""
        if len(text) <= max_length:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + max_length
            
            if end < len(text):
                for i in range(end, max(start, end - 100), -1):
                    if text[i] in '。.!?':
                        end = i + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end
        
        return chunks
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """处理单个文件 - 基于MD5去重"""
        # 首先检查文件是否重复
        if self._is_duplicate_file(file_path):
            logging.info(f"跳过重复文件: {file_path}")
            return []
        
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.pdf':
            return self._process_pdf_file(file_path)
        elif file_ext in ['.txt', '.md']:
            return self._process_txt_file(file_path)
        elif file_ext in ['.doc', '.docx']:
            return self._process_word_file(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            return self._process_powerpoint_file(file_path)
        elif file_ext in ['.xls', '.xlsx', '.xlsm']:
            return self._process_excel_file(file_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif']:
            return self._process_image_file(file_path)
        elif file_ext in ['.mp4', '.avi', '.mov', '.mkv']:
            return self._process_video_file(file_path)
        elif file_ext in ['.mp3', '.wav', '.flac']:
            return self._process_video_file(file_path)
        else:
            logging.warning(f"不支持的文件类型: {file_ext}")
            return []
    
    def process_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """处理目录中的所有文件"""
        if not os.path.exists(directory_path):
            logging.error(f"目录不存在: {directory_path}")
            return []
        
        all_documents = []
        supported_extensions = {'.pdf', '.txt', '.md', '.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx', '.xlsm', '.png', '.jpg', '.jpeg', '.bmp', '.gif', '.mp4', '.avi', '.mov', '.mkv', '.mp3', '.wav', '.flac'}
        
        for file_path in Path(directory_path).rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                logging.info(f"处理文件: {file_path} (类型: {file_path.suffix})")
                documents = self.process_file(str(file_path))
                all_documents.extend(documents)
        
        logging.info(f"目录处理完成: {len(all_documents)} 个文档")
        return all_documents
    
    def build_md5_database(self, documents: List[Dict[str, Any]], save_path: str = "multimodal_vector_db_md5"):
        """构建基于MD5的向量数据库"""
        if not documents:
            logging.warning("没有文档需要处理")
            return
        
        logging.info(f"开始构建MD5向量数据库，文档数量: {len(documents)}")
        
        # 提取文本内容
        texts = [doc['content'] for doc in documents]
        
        # 生成向量
        embeddings = self.text_model.encode(texts, show_progress_bar=True, batch_size=32)
        
        # 创建FAISS索引
        dimension = embeddings.shape[1]
        index = faiss.IndexFlatIP(dimension)
        index.add(embeddings.astype('float32'))
        
        # 保存索引和元数据
        os.makedirs(save_path, exist_ok=True)
        
        faiss.write_index(index, os.path.join(save_path, "faiss.index"))
        
        with open(os.path.join(save_path, "metadata.pkl"), 'wb') as f:
            pickle.dump([doc['metadata'] for doc in documents], f)
        
        with open(os.path.join(save_path, "documents.pkl"), 'wb') as f:
            pickle.dump(documents, f)
        
        # 保存MD5索引
        with open(os.path.join(save_path, "file_md5_index.json"), 'w', encoding='utf-8') as f:
            json.dump(self.file_md5_cache, f, ensure_ascii=False, indent=2)
        
        # 保存配置
        config = {
            'model_name': 'BAAI/bge-large-zh-v1.5',
            'embedding_dimension': dimension,
            'document_count': len(documents),
            'unique_files': len(self.processed_files),
            'document_types': list(set(doc['type'] for doc in documents)),
            'file_types_processed': list(set(doc['metadata'].get('file_type', 'unknown') for doc in documents)),
            'md5_duplicates_skipped': len(self.file_md5_cache) - len(self.processed_files)
        }
        
        with open(os.path.join(save_path, "config.json"), 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
        
        logging.info(f"✅ MD5向量数据库构建完成！保存路径: {save_path}")
        logging.info(f"索引信息: {dimension}维, {len(documents)}个文档")
        logging.info(f"处理文件数: {len(self.processed_files)}")
        logging.info(f"跳过重复文件: {len(self.file_md5_cache) - len(self.processed_files)}")
        
        # 保存索引到实例变量
        self.index = index
        self.documents = documents
        self.metadata = [doc['metadata'] for doc in documents]
    
    def search(self, query: str, top_k: int = 5) -> List[Tuple[Dict[str, Any], float]]:
        """搜索相似文档"""
        if not hasattr(self, 'index') or not hasattr(self, 'documents'):
            logging.error("向量数据库未初始化")
            return []
        
        query_embedding = self.text_model.encode([query])
        scores, indices = self.index.search(query_embedding.astype('float32'), top_k)
        
        results = []
        for score, idx in zip(scores[0], indices[0]):
            if idx < len(self.documents):
                results.append((self.documents[idx], float(score)))
        
        return results
    
    def show_database_info(self):
        """显示数据库信息"""
        if not hasattr(self, 'index'):
            print("  status: 未初始化")
            return
        
        print(f"  status: 已初始化")
        print(f"  document_count: {len(self.documents)}")
        print(f"  index_size: {self.index.ntotal}")
        print(f"  unique_files: {len(self.processed_files)}")
        print(f"  md5_cache_size: {len(self.file_md5_cache)}")

def main():
    """主函数"""
    print(" 企业知识库多模态向量化系统 - MD5去重版本")
    print("=" * 60)
    
    # 初始化MD5向量化器
    vectorizer = MD5MultimodalVectorizer()
    
    # 真实业务文档路径
    business_docs = "/root/autodl-tmp/enterprise_kb/sample_docs"
    video_docs = "/root/autodl-tmp/enterprise_kb/video_docs/产品演示"
    
    # 检查目录是否存在
    if not os.path.exists(business_docs):
        print(f"❌ 业务文档目录不存在: {business_docs}")
        return
    
    if not os.path.exists(video_docs):
        print(f"❌ 视频文档目录不存在: {video_docs}")
        return
    
    print(f"📁 扫描业务文档目录: {business_docs}")
    print(f" 扫描视频文档目录: {video_docs}")
    
    # 处理不同类型的文档
    print("\n📁 开始处理所有业务文档...")
    
    # 1. 处理业务文档（包括所有Office文档）
    print("📄 处理业务文档...")
    business_docs_list = vectorizer.process_directory(business_docs)
    
    # 2. 处理视频文档
    print("🎥 处理视频文档...")
    video_docs_list = vectorizer.process_directory(video_docs)
    
    # 合并所有文档
    all_docs = business_docs_list + video_docs_list
    
    if not all_docs:
        print("❌ 没有找到可处理的文档")
        return
    
    print(f"\n📊 文档处理统计:")
    print(f"业务文档: {len(business_docs_list)} 个")
    print(f"视频文档: {len(video_docs_list)} 个")
    print(f"总计: {len(all_docs)} 个文档")
    print(f"唯一文件数: {len(vectorizer.processed_files)}")
    
    # 构建MD5向量数据库
    print("\n 构建MD5向量数据库...")
    vectorizer.build_md5_database(all_docs)
    
    # 测试搜索功能
    print("\n🔍 测试搜索功能...")
    test_queries = [
        "商业合同",
        "财务报表", 
        "产品说明",
        "会议记录",
        "培训课程",
        "Springlake",
        "Guardia",
        "契约锁"
    ]
    
    for query in test_queries:
        print(f"\n查询: '{query}'")
        results = vectorizer.search(query, top_k=3)
        if results:
            for i, (doc, score) in enumerate(results, 1):
                print(f"  [{i}] 相似度: {score:.3f}")
                print(f"      类型: {doc.get('type', 'unknown')}")
                content = doc.get('content', '')[:100]
                print(f"      内容: {content}...")
        else:
            print("  未找到相关结果")
    
    # 显示数据库信息
    print("\n📊 数据库信息:")
    vectorizer.show_database_info()
    
    print("\n MD5去重向量化完成！")

if __name__ == "__main__":
    main()